"""Parallel processing for faster conversions."""

import os
import sys
import tempfile
from pathlib import Path
from typing import List
from concurrent.futures import ProcessPoolExecutor, ThreadPoolExecutor, as_completed


def convert_single_slide_playwright_pdf(html_file: Path, output_dir: str, index: int) -> tuple:
    """Convert a single HTML slide to PDF using Playwright."""
    try:
        from playwright.sync_api import sync_playwright
        
        with sync_playwright() as p:
            browser = p.chromium.launch()
            page = browser.new_page(viewport={'width': 1280, 'height': 720})
            
            page.goto(f"file://{html_file.absolute()}", timeout=30000)
            page.wait_for_load_state('networkidle', timeout=30000)
            
            # Create temp PDF
            temp_pdf = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf', dir=output_dir)
            page.pdf(path=temp_pdf.name, width='1280px', height='720px', print_background=True)
            
            browser.close()
            
            return (index, temp_pdf.name, None)
    except Exception as e:
        return (index, None, str(e))


def convert_single_slide_playwright_ppt(html_file: Path, output_dir: str, index: int) -> tuple:
    """Convert a single HTML slide to image using Playwright."""
    try:
        from playwright.sync_api import sync_playwright
        from PIL import Image
        import io
        
        with sync_playwright() as p:
            browser = p.chromium.launch()
            page = browser.new_page(viewport={'width': 1280, 'height': 720})
            
            page.goto(f"file://{html_file.absolute()}", timeout=30000)
            page.wait_for_load_state('networkidle', timeout=30000)
            
            # Take screenshot
            screenshot = page.screenshot()
            
            # Save to temp file
            img = Image.open(io.BytesIO(screenshot))
            temp_img = tempfile.NamedTemporaryFile(delete=False, suffix='.png', dir=output_dir)
            img.save(temp_img.name)
            
            browser.close()
            
            return (index, temp_img.name, None)
    except Exception as e:
        return (index, None, str(e))


def parallel_convert_to_pdf_playwright(html_files: List[Path], output_path: str, workers: int = 4, quiet: bool = False):
    """Convert HTML slides to PDF using parallel processing."""
    from PyPDF2 import PdfMerger
    
    if not quiet:
        print(f"Converting {len(html_files)} HTML slides to PDF using Playwright (parallel mode with {workers} workers)...")
    
    temp_dir = tempfile.mkdtemp()
    results = {}
    
    try:
        # Use ThreadPoolExecutor for I/O bound tasks (Playwright)
        with ThreadPoolExecutor(max_workers=workers) as executor:
            # Submit all tasks
            futures = {
                executor.submit(convert_single_slide_playwright_pdf, html_file, temp_dir, i): i
                for i, html_file in enumerate(html_files)
            }
            
            # Show progress
            try:
                from tqdm import tqdm
                progress = tqdm(total=len(html_files), desc="Converting", disable=quiet)
            except ImportError:
                progress = None
            
            # Collect results
            for future in as_completed(futures):
                index, pdf_path, error = future.result()
                
                if error:
                    if not quiet:
                        print(f"  Warning: Failed to process slide {index + 1}: {error}")
                else:
                    results[index] = pdf_path
                
                if progress:
                    progress.update(1)
            
            if progress:
                progress.close()
        
        if not results:
            print("Error: No slides were successfully converted")
            sys.exit(1)
        
        # Merge PDFs in order
        if not quiet:
            print("Merging PDFs...")
        
        merger = PdfMerger()
        for i in sorted(results.keys()):
            try:
                merger.append(results[i])
            except Exception as e:
                if not quiet:
                    print(f"  Warning: Failed to merge slide {i + 1}: {e}")
        
        merger.write(output_path)
        merger.close()
        
        if not quiet:
            print(f"✓ PDF created successfully: {output_path}")
            print(f"  ({len(results)}/{len(html_files)} slides converted)")
        
    except Exception as e:
        print(f"Error: {e}")
        sys.exit(1)
    finally:
        # Clean up temp files
        for pdf_path in results.values():
            try:
                if os.path.exists(pdf_path):
                    os.remove(pdf_path)
            except Exception:
                pass
        
        try:
            os.rmdir(temp_dir)
        except Exception:
            pass


def parallel_convert_to_ppt_playwright(html_files: List[Path], output_path: str, workers: int = 4, quiet: bool = False):
    """Convert HTML slides to PowerPoint using parallel processing."""
    from pptx import Presentation
    from pptx.util import Inches
    
    if not quiet:
        print(f"Converting {len(html_files)} HTML slides to PowerPoint using Playwright (parallel mode with {workers} workers)...")
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    temp_dir = tempfile.mkdtemp()
    results = {}
    
    try:
        # Use ThreadPoolExecutor for I/O bound tasks
        with ThreadPoolExecutor(max_workers=workers) as executor:
            # Submit all tasks
            futures = {
                executor.submit(convert_single_slide_playwright_ppt, html_file, temp_dir, i): i
                for i, html_file in enumerate(html_files)
            }
            
            # Show progress
            try:
                from tqdm import tqdm
                progress = tqdm(total=len(html_files), desc="Converting", disable=quiet)
            except ImportError:
                progress = None
            
            # Collect results
            for future in as_completed(futures):
                index, img_path, error = future.result()
                
                if error:
                    if not quiet:
                        print(f"  Warning: Failed to process slide {index + 1}: {error}")
                else:
                    results[index] = img_path
                
                if progress:
                    progress.update(1)
            
            if progress:
                progress.close()
        
        if not results:
            print("Error: No slides were successfully converted")
            sys.exit(1)
        
        # Add slides to presentation in order
        if not quiet:
            print("Creating PowerPoint...")
        
        for i in sorted(results.keys()):
            try:
                blank_slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_slide_layout)
                slide.shapes.add_picture(results[i], 0, 0, width=prs.slide_width, height=prs.slide_height)
            except Exception as e:
                if not quiet:
                    print(f"  Warning: Failed to add slide {i + 1}: {e}")
        
        # Save presentation
        prs.save(output_path)
        
        if not quiet:
            print(f"✓ PowerPoint created successfully: {output_path}")
            print(f"  ({len(results)}/{len(html_files)} slides converted)")
        
    except Exception as e:
        print(f"Error: {e}")
        sys.exit(1)
    finally:
        # Clean up temp files
        for img_path in results.values():
            try:
                if os.path.exists(img_path):
                    os.remove(img_path)
            except Exception:
                pass
        
        try:
            os.rmdir(temp_dir)
        except Exception:
            pass
