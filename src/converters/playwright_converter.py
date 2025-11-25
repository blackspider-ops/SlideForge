"""Playwright-based conversion methods."""

import os
import sys
import tempfile
from pathlib import Path
from typing import List


def convert_to_pdf_playwright(html_files: List[Path], output_path: str, quiet: bool = False, verbose: bool = False):
    """Convert HTML slides to PDF using Playwright."""
    from playwright.sync_api import sync_playwright
    from PyPDF2 import PdfMerger
    
    try:
        from tqdm import tqdm
        use_progress = not quiet
    except ImportError:
        use_progress = False
    
    if not quiet:
        print(f"Converting {len(html_files)} HTML slides to PDF using Playwright...")
    
    browser = None
    pdf_files = []
    
    try:
        with sync_playwright() as p:
            browser = p.chromium.launch()
            page = browser.new_page(viewport={'width': 1280, 'height': 720})
            
            for i, html_file in enumerate(html_files, 1):
                try:
                    print(f"  Processing slide {i}/{len(html_files)}: {html_file.name}")
                    
                    # Validate HTML file exists and is readable
                    if not html_file.exists():
                        print(f"  Warning: File not found, skipping: {html_file}")
                        continue
                    
                    page.goto(f"file://{html_file.absolute()}", timeout=30000)
                    page.wait_for_load_state('networkidle', timeout=30000)
                    
                    # Wait for fonts and icons to load
                    page.wait_for_timeout(2000)
                    
                    # Wait for document to be ready
                    page.evaluate("() => document.fonts.ready")
                    
                    # Force content to fit in one page by setting max-height
                    page.evaluate("""() => {
                        const slide = document.querySelector('.slide');
                        if (slide) {
                            slide.style.height = '720px';
                            slide.style.maxHeight = '720px';
                            slide.style.overflow = 'hidden';
                        }
                        document.body.style.height = '720px';
                        document.body.style.overflow = 'hidden';
                    }""")
                    
                    # Create temp PDF for this slide
                    temp_pdf = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
                    page.pdf(
                        path=temp_pdf.name, 
                        width='1280px', 
                        height='720px', 
                        print_background=True,
                        prefer_css_page_size=False
                    )
                    pdf_files.append(temp_pdf.name)
                    
                except Exception as e:
                    print(f"  Warning: Failed to process {html_file.name}: {e}")
                    continue
            
            browser.close()
        
        if not pdf_files:
            print("Error: No slides were successfully converted")
            sys.exit(1)
        
        # Merge all PDFs
        merger = PdfMerger()
        for pdf_file in pdf_files:
            try:
                merger.append(pdf_file)
            except Exception as e:
                print(f"  Warning: Failed to merge {pdf_file}: {e}")
        
        merger.write(output_path)
        merger.close()
        
        print(f"✓ PDF created successfully: {output_path}")
        
    except Exception as e:
        print(f"Error: {e}")
        print("Make sure you've run: playwright install chromium")
        sys.exit(1)
    finally:
        # Clean up temp files
        for pdf_file in pdf_files:
            try:
                if os.path.exists(pdf_file):
                    os.remove(pdf_file)
            except Exception:
                pass


def convert_to_ppt_playwright(html_files: List[Path], output_path: str):
    """Convert HTML slides to PowerPoint using Playwright."""
    from playwright.sync_api import sync_playwright
    from pptx import Presentation
    from pptx.util import Inches
    from PIL import Image
    import io
    
    print(f"Converting {len(html_files)} HTML slides to PowerPoint using Playwright...")
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    browser = None
    temp_images = []
    slides_created = 0
    
    try:
        with sync_playwright() as p:
            browser = p.chromium.launch()
            page = browser.new_page(viewport={'width': 1280, 'height': 720})
            
            for i, html_file in enumerate(html_files, 1):
                temp_img_path = None
                try:
                    print(f"  Processing slide {i}/{len(html_files)}: {html_file.name}")
                    
                    # Validate HTML file exists and is readable
                    if not html_file.exists():
                        print(f"  Warning: File not found, skipping: {html_file}")
                        continue
                    
                    page.goto(f"file://{html_file.absolute()}", timeout=30000)
                    page.wait_for_load_state('networkidle', timeout=30000)
                    
                    # Wait a bit for fonts to load
                    page.wait_for_timeout(500)
                    
                    # Take screenshot
                    screenshot = page.screenshot()
                    
                    # Add blank slide
                    blank_slide_layout = prs.slide_layouts[6]
                    slide = prs.slides.add_slide(blank_slide_layout)
                    
                    # Save screenshot to temp file
                    img = Image.open(io.BytesIO(screenshot))
                    temp_img = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
                    temp_img_path = temp_img.name
                    temp_images.append(temp_img_path)
                    img.save(temp_img_path)
                    
                    # Add image to slide
                    slide.shapes.add_picture(temp_img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
                    slides_created += 1
                    
                except Exception as e:
                    print(f"  Warning: Failed to process {html_file.name}: {e}")
                    continue
            
            browser.close()
        
        if slides_created == 0:
            print("Error: No slides were successfully converted")
            sys.exit(1)
        
        # Save presentation
        prs.save(output_path)
        print(f"✓ PowerPoint created successfully: {output_path}")
        print(f"  ({slides_created}/{len(html_files)} slides converted)")
        
    except Exception as e:
        print(f"Error: {e}")
        print("Make sure you've run: playwright install chromium")
        sys.exit(1)
    finally:
        # Clean up temp files
        for temp_img in temp_images:
            try:
                if os.path.exists(temp_img):
                    os.remove(temp_img)
            except Exception:
                pass
