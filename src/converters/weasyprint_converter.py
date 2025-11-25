"""WeasyPrint-based conversion methods."""

import os
import sys
import tempfile
from pathlib import Path
from typing import List


def convert_to_pdf_weasyprint(html_files: List[Path], output_path: str):
    """Convert HTML slides to PDF using WeasyPrint."""
    from weasyprint import HTML, CSS
    from PyPDF2 import PdfMerger
    
    print(f"Converting {len(html_files)} HTML slides to PDF using WeasyPrint...")
    
    pdf_files = []
    
    try:
        for i, html_file in enumerate(html_files, 1):
            try:
                print(f"  Processing slide {i}/{len(html_files)}: {html_file.name}")
                
                # Validate HTML file exists and is readable
                if not html_file.exists():
                    print(f"  Warning: File not found, skipping: {html_file}")
                    continue
                
                # Create temp PDF for this slide
                temp_pdf = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
                
                # Custom CSS for page size
                css = CSS(string='''
                    @page {
                        size: 1280px 720px;
                        margin: 0;
                    }
                ''')
                
                HTML(filename=str(html_file)).write_pdf(temp_pdf.name, stylesheets=[css])
                pdf_files.append(temp_pdf.name)
                
            except Exception as e:
                print(f"  Warning: Failed to process {html_file.name}: {e}")
                continue
        
        if not pdf_files:
            print("Error: No slides were successfully converted")
            sys.exit(1)
        
        # Merge all PDFs
        merger = PdfMerger()
        for pdf_file in pdf_files:
            try:
                merger.append(pdf_file)
            except Exception as e:
                print(f"  Warning: Failed to merge {pdf_file}: {e}")
        
        merger.write(output_path)
        merger.close()
        
        print(f"✓ PDF created successfully: {output_path}")
        
    except Exception as e:
        print(f"Error creating PDF: {e}")
        sys.exit(1)
    finally:
        # Clean up temp files
        for pdf_file in pdf_files:
            try:
                if os.path.exists(pdf_file):
                    os.remove(pdf_file)
            except Exception:
                pass


def convert_to_ppt_weasyprint(html_files: List[Path], output_path: str):
    """Convert HTML to PPT via WeasyPrint PDF then to images."""
    from weasyprint import HTML, CSS
    from pdf2image import convert_from_path
    from pptx import Presentation
    from pptx.util import Inches
    
    print(f"Converting {len(html_files)} HTML slides to PowerPoint using WeasyPrint...")
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    temp_files = []
    slides_created = 0
    
    try:
        for i, html_file in enumerate(html_files, 1):
            temp_pdf_path = None
            temp_img_path = None
            
            try:
                print(f"  Processing slide {i}/{len(html_files)}: {html_file.name}")
                
                # Validate HTML file exists and is readable
                if not html_file.exists():
                    print(f"  Warning: File not found, skipping: {html_file}")
                    continue
                
                # Convert HTML to PDF
                temp_pdf = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
                temp_pdf_path = temp_pdf.name
                temp_files.append(temp_pdf_path)
                
                css = CSS(string='@page { size: 1280px 720px; margin: 0; }')
                HTML(filename=str(html_file)).write_pdf(temp_pdf_path, stylesheets=[css])
                
                # Convert PDF to image
                images = convert_from_path(temp_pdf_path, dpi=150)
                
                if not images:
                    print(f"  Warning: No image generated for {html_file.name}")
                    continue
                
                # Save image
                temp_img = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
                temp_img_path = temp_img.name
                temp_files.append(temp_img_path)
                images[0].save(temp_img_path, 'PNG')
                
                # Add to presentation
                blank_slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_slide_layout)
                slide.shapes.add_picture(temp_img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
                slides_created += 1
                
            except Exception as e:
                print(f"  Warning: Failed to process {html_file.name}: {e}")
                continue
        
        if slides_created == 0:
            print("Error: No slides were successfully converted")
            sys.exit(1)
        
        prs.save(output_path)
        print(f"✓ PowerPoint created successfully: {output_path}")
        print(f"  ({slides_created}/{len(html_files)} slides converted)")
        
    except Exception as e:
        print(f"Error: {e}")
        if "poppler" in str(e).lower():
            print("\nNote: WeasyPrint PPT conversion requires poppler")
            print("Install with: brew install poppler (macOS)")
        sys.exit(1)
    finally:
        # Clean up temp files
        for temp_file in temp_files:
            try:
                if os.path.exists(temp_file):
                    os.remove(temp_file)
            except Exception:
                pass
