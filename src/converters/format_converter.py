"""Convert between PDF and PPT formats."""

import sys
from pathlib import Path


def convert_pdf_to_ppt(pdf_path: str, output_path: str, quiet: bool = False):
    """Convert PDF to PowerPoint."""
    try:
        from pdf2image import convert_from_path
        from pptx import Presentation
        from pptx.util import Inches
        import tempfile
        import os
    except ImportError:
        print("Error: Missing dependencies for PDF to PPT conversion")
        print("Install with: pip install pdf2image python-pptx")
        print("Note: Also requires poppler (brew install poppler on macOS)")
        sys.exit(1)
    
    if not quiet:
        print(f"Converting PDF to PowerPoint: {pdf_path}")
    
    try:
        # Convert PDF pages to images
        if not quiet:
            print("  Converting PDF pages to images...")
        images = convert_from_path(pdf_path, dpi=150)
        
        if not images:
            print("Error: No pages found in PDF")
            sys.exit(1)
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)
        
        temp_files = []
        
        if not quiet:
            print(f"  Creating PowerPoint with {len(images)} slides...")
        
        for i, image in enumerate(images, 1):
            if not quiet:
                print(f"    Processing page {i}/{len(images)}")
            
            # Save image to temp file
            temp_img = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
            temp_files.append(temp_img.name)
            image.save(temp_img.name, 'PNG')
            
            # Add slide
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(temp_img.name, 0, 0, width=prs.slide_width, height=prs.slide_height)
        
        # Save presentation
        prs.save(output_path)
        
        # Clean up temp files
        for temp_file in temp_files:
            try:
                os.remove(temp_file)
            except Exception:
                pass
        
        if not quiet:
            print(f"✓ PowerPoint created: {output_path}")
            print(f"  ({len(images)} pages converted)")
        
    except Exception as e:
        print(f"Error converting PDF to PPT: {e}")
        if "poppler" in str(e).lower():
            print("\nNote: PDF to PPT conversion requires poppler")
            print("Install with: brew install poppler (macOS)")
        sys.exit(1)


def convert_ppt_to_pdf(ppt_path: str, output_path: str, quiet: bool = False):
    """Convert PowerPoint to PDF."""
    try:
        from pptx import Presentation
    except ImportError:
        print("Error: python-pptx not installed")
        print("Install with: pip install python-pptx")
        sys.exit(1)
    
    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import landscape, letter
    except ImportError:
        print("Error: reportlab not installed")
        print("Install with: pip install reportlab")
        sys.exit(1)
    
    try:
        from PIL import Image
        import tempfile
        import os
    except ImportError:
        print("Error: pillow not installed")
        print("Install with: pip install pillow")
        sys.exit(1)
    
    if not quiet:
        print(f"Converting PowerPoint to PDF: {ppt_path}")
    
    try:
        # Load presentation
        prs = Presentation(ppt_path)
        
        if not prs.slides:
            print("Error: No slides found in PowerPoint")
            sys.exit(1)
        
        if not quiet:
            print(f"  Found {len(prs.slides)} slides")
            print("  Note: This creates a basic PDF. For best results, use 'Export as PDF' in PowerPoint.")
        
        # Create PDF
        # Note: This is a simplified conversion
        # For production use, consider using LibreOffice or PowerPoint's native export
        
        from PyPDF2 import PdfMerger
        import subprocess
        
        # Try using LibreOffice if available (best quality)
        try:
            if not quiet:
                print("  Attempting conversion with LibreOffice...")
            
            result = subprocess.run(
                ['soffice', '--headless', '--convert-to', 'pdf', '--outdir', 
                 str(Path(output_path).parent), ppt_path],
                capture_output=True,
                timeout=60
            )
            
            if result.returncode == 0:
                # LibreOffice creates filename.pdf in output dir
                expected_pdf = Path(output_path).parent / (Path(ppt_path).stem + '.pdf')
                if expected_pdf.exists():
                    if str(expected_pdf) != output_path:
                        expected_pdf.rename(output_path)
                    
                    if not quiet:
                        print(f"✓ PDF created: {output_path}")
                    return
            else:
                if not quiet:
                    print(f"  LibreOffice conversion failed: {result.stderr.decode()}")
        except FileNotFoundError:
            if not quiet:
                print("  LibreOffice not found")
        except subprocess.TimeoutExpired:
            if not quiet:
                print("  LibreOffice conversion timed out")
        
        # Fallback: Extract images from PPT and create PDF
        print("\n⚠️  LibreOffice not available - using image extraction method")
        print("For best results, install LibreOffice: brew install libreoffice (macOS)\n")
        
        # Try to extract images from slides
        from reportlab.pdfgen import canvas
        from reportlab.lib.utils import ImageReader
        import io
        
        if not quiet:
            print("  Extracting slide images...")
        
        # Create PDF with reportlab
        c = canvas.Canvas(output_path, pagesize=(1280, 720))
        
        slides_processed = 0
        for i, slide in enumerate(prs.slides, 1):
            if not quiet:
                print(f"    Processing slide {i}/{len(prs.slides)}")
            
            # Try to extract images from slide (slides created from screenshots have images)
            image_found = False
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Picture type
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        img_reader = ImageReader(io.BytesIO(image_bytes))
                        
                        # Draw image to fill the page
                        c.drawImage(img_reader, 0, 0, width=1280, height=720, preserveAspectRatio=False)
                        image_found = True
                        slides_processed += 1
                        break
                    except Exception as e:
                        if not quiet:
                            print(f"      Warning: Could not extract image from slide {i}: {e}")
            
            # If no image found, try to render text content
            if not image_found:
                y_position = 650
                c.setFont("Helvetica-Bold", 24)
                c.drawString(50, y_position, f"Slide {i}")
                y_position -= 40
                
                c.setFont("Helvetica", 14)
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # Wrap text
                        text = shape.text.strip()
                        words = text.split()
                        line = ""
                        for word in words:
                            test_line = line + word + " "
                            if len(test_line) > 80:
                                c.drawString(50, y_position, line)
                                y_position -= 20
                                line = word + " "
                            else:
                                line = test_line
                        if line:
                            c.drawString(50, y_position, line)
                            y_position -= 30
                        
                        if y_position < 50:
                            break
                
                slides_processed += 1
            
            c.showPage()
        
        c.save()
        
        if not quiet:
            if slides_processed > 0:
                print(f"✓ PDF created: {output_path}")
                print(f"  ({slides_processed}/{len(prs.slides)} slides converted)")
            else:
                print(f"⚠️  PDF created but may be empty: {output_path}")
            print("  (For better quality, install LibreOffice: brew install libreoffice)")
        
    except Exception as e:
        print(f"Error converting PPT to PDF: {e}")
        sys.exit(1)
